#!/usr/bin/env python3
"""
Insert 3 new technical capability slides into the existing Simon Meeting PPTX.
Slides 16, 17, 18 are inserted BEFORE the current closing slide (slide 15).
After insertion the deck has 18 slides total.

Strategy:
  1. Build the 3 new slides by appending to the presentation (they land at positions 15–17).
  2. Move slide 15 (the closing slide) to position 18 using lxml manipulation.
  3. Save back to the same file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import copy, os

# ── Theme (identical to generate_simon_meeting_pptx.py) ─────────────────────
BG_COLOR    = RGBColor(0x0F, 0x14, 0x19)
BG_LIGHTER  = RGBColor(0x1A, 0x1F, 0x26)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0xD4, 0xAA)
LIGHT_GRAY  = RGBColor(0xB0, 0xB8, 0xC4)
MID_GRAY    = RGBColor(0x6B, 0x73, 0x80)
DARK_BORDER = RGBColor(0x2A, 0x30, 0x3A)

FONT_NAME    = "Calibri"
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TOTAL_SLIDES = 18   # final deck size


# ── Low-level helpers ─────────────────────────────────────────────────────────

def set_slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def top_teal_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def add_accent_bar(slide, left=Inches(0.8), top=Inches(1.33),
                   width=Inches(1.5), height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_title(slide, text, top=Inches(0.45), left=Inches(0.8),
              width=Inches(11.7), size=Pt(32)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME


def add_subtitle(slide, text, top=Inches(1.65), left=Inches(0.8),
                 width=Inches(11.7)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(15)
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = FONT_NAME


def add_slide_number(slide, num, total=TOTAL_SLIDES):
    txBox = slide.shapes.add_textbox(Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    run.font.size = Pt(10)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT_NAME


def add_footer(slide, text, top=Inches(6.6)):
    txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.7), Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT_NAME


def add_bold_label_run(p, label, size=Pt(16)):
    r = p.add_run()
    r.text = label
    r.font.size = size
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT_NAME


def add_plain_run(p, text, size=Pt(16), color=None):
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.color.rgb = color if color else LIGHT_GRAY
    r.font.name = FONT_NAME


# ── Slide builders ────────────────────────────────────────────────────────────

def build_slide_16(prs):
    """Computational Infrastructure — 11 NVIDIA NIMs Live."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    top_teal_bar(slide)
    add_title(slide, "What We Can Compute \u2014 11 NVIDIA NIMs Live", size=Pt(31))
    add_accent_bar(slide)

    nims = [
        ("DiffDock v2.2",        "Molecular docking (drug \u2192 protein binding prediction)"),
        ("AlphaFold2",           "Protein structure prediction from sequence"),
        ("AlphaFold2-Multimer",  "Protein complex structures (multi-chain)"),
        ("ESMfold",              "Fast protein folding (no MSA needed, seconds)"),
        ("RFdiffusion",          "De novo protein binder design"),
        ("ProteinMPNN",          "Sequence optimization for designed proteins"),
        ("GenMol",               "AI molecule generation for target pockets"),
        ("MolMIM",               "Molecule optimization with QED scoring"),
        ("MSA Search",           "ColabFold sequence alignment"),
        ("ESM-2",                "Protein embeddings (1280-dim structural fingerprints)"),
        ("RDKit",                "Drug-likeness, BBB, CNS MPO, PAINS screening"),
    ]

    # Two-column grid: 6 left, 5 right
    col_left  = nims[:6]
    col_right = nims[6:]

    col_configs = [
        (Inches(0.75),  col_left),
        (Inches(7.05),  col_right),
    ]

    item_h  = Inches(0.62)
    y_start = Inches(1.75)

    for col_left_x, items in col_configs:
        for i, (name, desc) in enumerate(items):
            y = y_start + i * item_h

            # subtle card strip
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                col_left_x, y, Inches(6.0), item_h - Inches(0.06)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = BG_LIGHTER
            card.line.color.rgb = DARK_BORDER
            card.line.width = Pt(0.75)

            # NIM name (teal, bold)
            txN = slide.shapes.add_textbox(
                col_left_x + Inches(0.18), y + Inches(0.05),
                Inches(1.9), item_h - Inches(0.1)
            )
            tfN = txN.text_frame
            pN  = tfN.paragraphs[0]
            rN  = pN.add_run()
            rN.text = name
            rN.font.size = Pt(13)
            rN.font.bold = True
            rN.font.color.rgb = TEAL
            rN.font.name = FONT_NAME

            # Em dash separator
            txSep = slide.shapes.add_textbox(
                col_left_x + Inches(2.1), y + Inches(0.05),
                Inches(0.2), item_h - Inches(0.1)
            )
            tfSep = txSep.text_frame
            pSep  = tfSep.paragraphs[0]
            rSep  = pSep.add_run()
            rSep.text = "\u2014"
            rSep.font.size = Pt(13)
            rSep.font.color.rgb = MID_GRAY
            rSep.font.name = FONT_NAME

            # Description (light gray)
            txD = slide.shapes.add_textbox(
                col_left_x + Inches(2.35), y + Inches(0.05),
                Inches(3.5), item_h - Inches(0.1)
            )
            tfD = txD.text_frame
            tfD.word_wrap = True
            pD  = tfD.paragraphs[0]
            rD  = pD.add_run()
            rD.text = desc
            rD.font.size = Pt(12)
            rD.font.color.rgb = LIGHT_GRAY
            rD.font.name = FONT_NAME

    add_footer(slide, "All accessible via REST API + MCP Server for AI agent integration")
    add_slide_number(slide, 16)
    return slide


def build_slide_17(prs):
    """What We Can Do For You — Computational Services."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    top_teal_bar(slide)
    add_title(slide, "Computational Services We Offer", size=Pt(32))
    add_accent_bar(slide)

    services = [
        ("Target Analysis",
         "Score any gene/protein across 5 evidence dimensions from 14,176 claims"),
        ("Virtual Screening",
         "Dock your compound library against any SMA target (DiffDock v2.2)"),
        ("Protein Binder Design",
         "Design novel proteins that bind ROCK2/LIMK1/MAPK14 (RFdiffusion + ProteinMPNN)"),
        ("Drug Discovery",
         "Generate + filter + dock novel molecules for your target (GenMol + RDKit + DiffDock)"),
        ("Literature Synthesis",
         "Cross-paper analysis across 6,535 sources — find connections no single paper reveals"),
        ("Hypothesis Generation",
         "AI-generated testable predictions with experiment proposals"),
    ]

    # 2-column, 3-row grid of capability cards
    card_w  = Inches(5.85)
    card_h  = Inches(1.35)
    gap_x   = Inches(0.45)
    gap_y   = Inches(0.28)
    y_start = Inches(1.8)
    left1   = Inches(0.65)
    left2   = Inches(0.65) + card_w + gap_x

    for i, (name, desc) in enumerate(services):
        col  = i % 2
        row  = i // 2
        left = left1 if col == 0 else left2
        top  = y_start + row * (card_h + gap_y)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_LIGHTER
        card.line.color.rgb = TEAL
        card.line.width = Pt(1.5)

        # Number badge
        num_str = str(i + 1)
        txNum = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.12), Inches(0.38), Inches(0.5)
        )
        tfNum = txNum.text_frame
        pNum  = tfNum.paragraphs[0]
        rNum  = pNum.add_run()
        rNum.text = num_str
        rNum.font.size = Pt(20)
        rNum.font.bold = True
        rNum.font.color.rgb = TEAL
        rNum.font.name = FONT_NAME

        # Service name
        txName = slide.shapes.add_textbox(
            left + Inches(0.65), top + Inches(0.1), card_w - Inches(0.8), Inches(0.52)
        )
        tfName = txName.text_frame
        tfName.word_wrap = True
        pName  = tfName.paragraphs[0]
        rName  = pName.add_run()
        rName.text = name
        rName.font.size = Pt(15)
        rName.font.bold = True
        rName.font.color.rgb = WHITE
        rName.font.name = FONT_NAME

        # Description
        txDesc = slide.shapes.add_textbox(
            left + Inches(0.65), top + Inches(0.66), card_w - Inches(0.8), Inches(0.6)
        )
        tfDesc = txDesc.text_frame
        tfDesc.word_wrap = True
        pDesc  = tfDesc.paragraphs[0]
        rDesc  = pDesc.add_run()
        rDesc.text = desc
        rDesc.font.size = Pt(12)
        rDesc.font.color.rgb = LIGHT_GRAY
        rDesc.font.name = FONT_NAME

    add_footer(slide, "All open source. Results delivered as structured data, ready for your analysis.")
    add_slide_number(slide, 17)
    return slide


def build_slide_18(prs):
    """Research Roadmap — Next 4 Weeks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    top_teal_bar(slide)
    add_title(slide, "What\u2019s Coming Next", size=Pt(34))
    add_accent_bar(slide)
    add_subtitle(slide, "Research Roadmap \u2014 Next 4 Weeks", top=Inches(1.65))

    weeks = [
        ("Week 1", "Cross-disease transcriptomics",
         "20 ALS datasets \u00D7 24-gene panel — find shared and divergent expression signatures"),
        ("Week 2", "Extended DiffDock campaign",
         "63 targets \u00D7 top drug candidates — full binding landscape across the SMA target space"),
        ("Week 3", "Protein binder design",
         "Novel proteins targeting ROCK2 / MAPK14 / LIMK1 (RFdiffusion + ProteinMPNN pipeline)"),
        ("Week 4", "Molecule generation + hypothesis mining",
         "GenMol generates novel scaffolds — filtered by RDKit, docked by DiffDock, ranked by evidence"),
    ]

    card_w  = Inches(10.9)
    card_h  = Inches(1.15)
    gap_y   = Inches(0.25)
    y_start = Inches(2.2)
    left    = Inches(1.2)

    for i, (week, title, desc) in enumerate(weeks):
        top = y_start + i * (card_h + gap_y)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_LIGHTER
        card.line.color.rgb = TEAL
        card.line.width = Pt(2)

        # Week badge (left, teal, large)
        txW = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.1), Inches(0.95), card_h - Inches(0.2)
        )
        tfW = txW.text_frame
        pW  = tfW.paragraphs[0]
        pW.alignment = PP_ALIGN.CENTER
        rW  = pW.add_run()
        rW.text = week
        rW.font.size = Pt(13)
        rW.font.bold = True
        rW.font.color.rgb = TEAL
        rW.font.name = FONT_NAME

        # Vertical separator line
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left + Inches(1.25), top + Inches(0.15),
            Inches(0.03), card_h - Inches(0.3)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = DARK_BORDER
        sep.line.fill.background()

        # Title + description stacked
        txTitle = slide.shapes.add_textbox(
            left + Inches(1.4), top + Inches(0.08), card_w - Inches(1.6), Inches(0.42)
        )
        tfTitle = txTitle.text_frame
        tfTitle.word_wrap = True
        pTitle  = tfTitle.paragraphs[0]
        rTitle  = pTitle.add_run()
        rTitle.text = title
        rTitle.font.size = Pt(15)
        rTitle.font.bold = True
        rTitle.font.color.rgb = WHITE
        rTitle.font.name = FONT_NAME

        txDesc = slide.shapes.add_textbox(
            left + Inches(1.4), top + Inches(0.55), card_w - Inches(1.6), Inches(0.52)
        )
        tfDesc = txDesc.text_frame
        tfDesc.word_wrap = True
        pDesc  = tfDesc.paragraphs[0]
        rDesc  = pDesc.add_run()
        rDesc.text = desc
        rDesc.font.size = Pt(12)
        rDesc.font.color.rgb = LIGHT_GRAY
        rDesc.font.name = FONT_NAME

    # Bottom summary banner
    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.2), Inches(6.25), Inches(10.9), Inches(0.58)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0x00, 0x2B, 0x22)   # dark teal tint
    banner.line.color.rgb = TEAL
    banner.line.width = Pt(1)

    txBanner = slide.shapes.add_textbox(
        Inches(1.4), Inches(6.28), Inches(10.5), Inches(0.5)
    )
    tfBanner = txBanner.text_frame
    tfBanner.word_wrap = True
    pBanner  = tfBanner.paragraphs[0]
    pBanner.alignment = PP_ALIGN.CENTER
    rBanner  = pBanner.add_run()
    rBanner.text = "All computation runs autonomously. Results feed back into the platform daily."
    rBanner.font.size = Pt(14)
    rBanner.font.color.rgb = TEAL
    rBanner.font.name = FONT_NAME

    add_slide_number(slide, 18)
    return slide


# ── Slide-reorder helper ──────────────────────────────────────────────────────

def move_slide(prs, old_index, new_index):
    """
    Move a slide from old_index to new_index (0-based) in the slide list.
    Uses lxml to manipulate the presentation's <p:sldIdLst>.
    """
    xml_slides = prs.slides._sldIdLst
    slides     = list(xml_slides)
    slide_elem = slides[old_index]
    xml_slides.remove(slide_elem)
    xml_slides.insert(new_index, slide_elem)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    pptx_path = (
        "/mnt/c/Users/bryza/Dropbox/Christian fischer/SMA/sma-platform/"
        "docs/presentations/SMA_Research_Platform_Simon_Meeting_March2026.pptx"
    )

    print(f"Opening: {pptx_path}")
    prs = Presentation(pptx_path)
    print(f"Slide count before: {len(prs.slides)}")

    if len(prs.slides) != 15:
        raise RuntimeError(
            f"Expected 15 slides, found {len(prs.slides)}. "
            "Aborting to avoid double-insertion."
        )

    # ── Step 1: Append 3 new slides (they land at indices 15, 16, 17) ──
    build_slide_16(prs)   # lands at index 15
    build_slide_17(prs)   # lands at index 16
    build_slide_18(prs)   # lands at index 17

    # Closing slide is currently at index 14 (slide 15).
    # New slides are at indices 15, 16, 17.
    # We need the final order: 1-14 (unchanged), 16-new, 17-new, 18-new, 15-closing.
    # That means moving the closing slide (index 14) to index 17 (last position).

    # ── Step 2: Move closing slide to the end ──
    move_slide(prs, old_index=14, new_index=17)

    print(f"Slide count after: {len(prs.slides)}")
    assert len(prs.slides) == 18, f"Expected 18 slides, got {len(prs.slides)}"

    # ── Step 3: Save back to same file ──
    prs.save(pptx_path)
    print(f"Saved: {pptx_path}")
    print("Done — 3 technical slides inserted before closing slide. Total: 18 slides.")


if __name__ == "__main__":
    main()
