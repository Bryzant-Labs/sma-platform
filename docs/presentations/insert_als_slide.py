#!/usr/bin/env python3
"""
Insert CORO1C-ALS Cross-Disease slide after slide 8 (Actin Pathway)
into SMA_Research_Platform_Simon_Meeting_March2026.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import copy
import os

# --- Theme colors ---
BG_COLOR    = RGBColor(0x0F, 0x14, 0x19)
BG_LIGHTER  = RGBColor(0x1A, 0x1F, 0x26)
BG_DARKCARD = RGBColor(0x16, 0x1B, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0xD4, 0xAA)
LIGHT_GRAY  = RGBColor(0xB0, 0xB8, 0xC4)
MID_GRAY    = RGBColor(0x6B, 0x73, 0x80)
DARK_BORDER = RGBColor(0x2A, 0x30, 0x3A)
GOLD        = RGBColor(0xFF, 0xCC, 0x44)   # Star highlight

FONT_NAME = "Calibri"
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_number(slide, num, total=19):
    txBox = slide.shapes.add_textbox(Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    run.font.size = Pt(10)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT_NAME


def add_accent_bar(slide, left=Inches(0.8), top=None, width=Inches(1.5), height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_textbox(slide, left, top, width, height, text, size, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_NAME
    return txBox


def make_als_cross_disease_slide(prs):
    """Build the CORO1C-ALS Cross-Disease Validation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    # ── Title ──────────────────────────────────────────────────────────────────
    add_textbox(
        slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.9),
        "Cross-Disease Validation: Actin Genes in ALS Transcriptomics",
        size=28, bold=True, color=WHITE
    )

    # Accent bar under title
    add_accent_bar(slide, top=Inches(1.15))

    # Subtitle
    add_textbox(
        slide, Inches(0.8), Inches(1.25), Inches(11.7), Inches(0.4),
        "GSE113924 — PFN1-G118V ALS Mouse Model (Barham et al. 2018, PMID 30213953)",
        size=13, bold=False, color=LIGHT_GRAY
    )

    # ── Table setup ────────────────────────────────────────────────────────────
    # Columns: Gene | ALS (GSE113924) | SMA Direction | Interpretation
    col_widths = [Inches(1.5), Inches(2.8), Inches(2.4), Inches(5.8)]
    col_lefts = [Inches(0.5)]
    for w in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + w)

    rows = [
        # (gene, als, sma, interp, highlight_teal)
        ("LIMK1",  "DOWN (padj 0.0005)", "Predicted DOWN", "Validates ROCK-LIMK axis \u2605", True),
        ("LIMK2",  "DOWN (padj 0.003)",  "Predicted DOWN", "Both LIMKs compromised",         True),
        ("CORO1C", "UP (padj 0.003)",    "UP (GSE69175)",  "Shared actin stress response",    True),
        ("PLS3",   "DOWN (padj 0.011)",  "Protective modifier", "SMA protector lost in ALS", True),
        ("Arp2/3", "All 5 UP",           "\u2014",         "Actin nucleation stress",         False),
        ("CFL2",   "NOT significant",    "UP 2.9x in SMA", "Possibly SMA-specific",           False),
    ]

    header_top = Inches(1.75)
    row_h      = Inches(0.58)

    # Header row
    header_cells = ["Gene", "ALS (GSE113924)", "SMA Direction", "Interpretation"]
    for ci, (left, text) in enumerate(zip(col_lefts, header_cells)):
        # Header background
        hdr_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, header_top,
            col_widths[ci], row_h - Inches(0.04)
        )
        hdr_bg.fill.solid()
        hdr_bg.fill.fore_color.rgb = RGBColor(0x00, 0x5C, 0x4A)  # dark teal
        hdr_bg.line.color.rgb = DARK_BORDER
        hdr_bg.line.width = Pt(0.75)

        txH = slide.shapes.add_textbox(
            left + Inches(0.08), header_top + Inches(0.1),
            col_widths[ci] - Inches(0.12), row_h - Inches(0.2)
        )
        tfH = txH.text_frame
        tfH.word_wrap = True
        pH = tfH.paragraphs[0]
        pH.alignment = PP_ALIGN.LEFT
        rH = pH.add_run()
        rH.text = text
        rH.font.size = Pt(13)
        rH.font.bold = True
        rH.font.color.rgb = TEAL
        rH.font.name = FONT_NAME

    # Data rows
    for ri, (gene, als, sma, interp, highlight) in enumerate(rows):
        row_top = header_top + row_h + ri * row_h
        row_color = BG_DARKCARD if ri % 2 == 0 else BG_LIGHTER
        cell_texts = [gene, als, sma, interp]

        for ci, (left, text) in enumerate(zip(col_lefts, cell_texts)):
            cell_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, row_top,
                col_widths[ci], row_h - Inches(0.04)
            )
            cell_bg.fill.solid()
            cell_bg.fill.fore_color.rgb = row_color
            cell_bg.line.color.rgb = DARK_BORDER
            cell_bg.line.width = Pt(0.5)

            txC = slide.shapes.add_textbox(
                left + Inches(0.08), row_top + Inches(0.09),
                col_widths[ci] - Inches(0.12), row_h - Inches(0.18)
            )
            tfC = txC.text_frame
            tfC.word_wrap = True
            pC = tfC.paragraphs[0]
            pC.alignment = PP_ALIGN.LEFT

            # Gene column: bold white; ALS/SMA/Interp: teal if highlight, else light_gray
            if ci == 0:
                rC = pC.add_run()
                rC.text = text
                rC.font.size = Pt(13)
                rC.font.bold = True
                rC.font.color.rgb = WHITE
                rC.font.name = FONT_NAME
            elif ci == 3 and highlight and "\u2605" in text:
                # Split star from text
                parts = text.split("\u2605")
                rC = pC.add_run()
                rC.text = parts[0].rstrip()
                rC.font.size = Pt(12)
                rC.font.color.rgb = TEAL if highlight else LIGHT_GRAY
                rC.font.name = FONT_NAME
                rStar = pC.add_run()
                rStar.text = " \u2605"
                rStar.font.size = Pt(13)
                rStar.font.bold = True
                rStar.font.color.rgb = GOLD
                rStar.font.name = FONT_NAME
            else:
                rC = pC.add_run()
                rC.text = text
                rC.font.size = Pt(12)
                rC.font.bold = (ci == 0)
                if highlight:
                    rC.font.color.rgb = TEAL if ci in (1, 2) else WHITE
                else:
                    rC.font.color.rgb = LIGHT_GRAY if ci in (1, 2) else MID_GRAY
                rC.font.name = FONT_NAME

    # ── Footer bar ─────────────────────────────────────────────────────────────
    footer_top = header_top + row_h + len(rows) * row_h + Inches(0.12)

    footer_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), footer_top,
        Inches(12.3), Inches(0.52)
    )
    footer_bg.fill.solid()
    footer_bg.fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x31)  # very dark teal
    footer_bg.line.fill.background()

    txF = slide.shapes.add_textbox(
        Inches(0.65), footer_top + Inches(0.06),
        Inches(12.0), Inches(0.42)
    )
    tfF = txF.text_frame
    tfF.word_wrap = True
    pF = tfF.paragraphs[0]
    pF.alignment = PP_ALIGN.LEFT
    rF1 = pF.add_run()
    rF1.text = "\u2605  "
    rF1.font.size = Pt(13)
    rF1.font.bold = True
    rF1.font.color.rgb = GOLD
    rF1.font.name = FONT_NAME
    rF2 = pF.add_run()
    rF2.text = ("LIMK1 downregulation is the strongest finding — "
                "validates our ROCK-LIMK-Cofilin hypothesis with independent ALS data")
    rF2.font.size = Pt(13)
    rF2.font.bold = True
    rF2.font.color.rgb = WHITE
    rF2.font.name = FONT_NAME

    # ── Caveat line (small, bottom) ─────────────────────────────────────────────
    caveat_top = footer_top + Inches(0.60)
    txCav = slide.shapes.add_textbox(
        Inches(0.5), caveat_top,
        Inches(12.3), Inches(0.32)
    )
    tfCav = txCav.text_frame
    tfCav.word_wrap = True
    pCav = tfCav.paragraphs[0]
    pCav.alignment = PP_ALIGN.LEFT
    rCav = pCav.add_run()
    rCav.text = (
        "Computational re-analysis of published data. Bulk RNA-seq. "
        "Requires validation in human ALS datasets (GSE153960, n=380)."
    )
    rCav.font.size = Pt(10)
    rCav.font.color.rgb = MID_GRAY
    rCav.font.name = FONT_NAME
    rCav.font.italic = True

    add_slide_number(slide, 9, total=19)
    return slide


def move_slide_to_position(prs, from_index, to_index):
    """Move slide from from_index to to_index (0-based) in the presentation."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    # Remove the slide element from its current position
    el = slides_list[from_index]
    xml_slides.remove(el)
    # Insert at new position
    xml_slides.insert(to_index, el)


PPTX_PATH = "/mnt/c/Users/bryza/Dropbox/Christian fischer/SMA/sma-platform/docs/presentations/SMA_Research_Platform_Simon_Meeting_March2026.pptx"
TEMP_PATH = "/tmp/SMA_Simon_Meeting_new.pptx"

prs = Presentation(PPTX_PATH)
print(f"Before: {len(prs.slides)} slides")

# Add the new slide (it goes to the end)
make_als_cross_disease_slide(prs)

print(f"After adding: {len(prs.slides)} slides")

# Move slide from last position (index 18) to position 8 (after slide 8, before old slide 9)
# New slide is at index 18 (0-based), needs to be at index 8
move_slide_to_position(prs, len(prs.slides._sldIdLst) - 1, 8)

print(f"After moving: {len(prs.slides)} slides")

# Verify order
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            first = shape.text_frame.paragraphs[0].text.strip()
            if first:
                print(f"  Slide {i+1}: {first[:60]}")
                break

# Save to temp first, then copy over original
prs.save(TEMP_PATH)
print(f"\nSaved to temp. Final slide count: {len(prs.slides)}")

import shutil
shutil.copy2(TEMP_PATH, PPTX_PATH)
print(f"Copied to: {PPTX_PATH}")
