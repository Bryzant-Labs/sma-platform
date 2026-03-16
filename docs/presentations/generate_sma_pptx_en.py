#!/usr/bin/env python3
"""
Generate SMA Research Platform presentation (English, dark theme).
Output: SMA_Research_Platform_EN.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Theme colors ---
BG_COLOR = RGBColor(0x0F, 0x14, 0x19)       # Dark navy/charcoal
BG_LIGHTER = RGBColor(0x1A, 0x1F, 0x26)     # Slightly lighter for cards
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x00, 0xD4, 0xAA)           # Accent
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC4)     # Secondary text
MID_GRAY = RGBColor(0x6B, 0x73, 0x80)       # Muted text
DARK_BORDER = RGBColor(0x2A, 0x30, 0x3A)    # Borders/dividers

FONT_NAME = "Calibri"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_COLOR):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_number(slide, num, total=21):
    """Add slide number in bottom-right corner."""
    txBox = slide.shapes.add_textbox(
        Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    run.font.size = Pt(10)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT_NAME


def add_accent_bar(slide, left=Inches(0.8), top=None, width=Inches(1.5), height=Inches(0.06)):
    """Add a teal accent bar under the title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    shape.shadow.inherit = False


def make_title_slide(prs, slide_num):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    # Decorative top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "SMA Research Platform"
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(3.6), Inches(9.3), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "AI-Driven Drug Discovery for Spinal Muscular Atrophy"
    run2.font.size = Pt(24)
    run2.font.color.rgb = TEAL
    run2.font.name = FONT_NAME

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(4.8), Inches(3.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = DARK_BORDER
    line.line.fill.background()

    # Footer
    txBox3 = slide.shapes.add_textbox(Inches(2.0), Inches(5.3), Inches(9.3), Inches(0.8))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "sma-research.info  |  Open Source  |  AGPL-3.0"
    run3.font.size = Pt(16)
    run3.font.color.rgb = MID_GRAY
    run3.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def make_bullet_slide(prs, slide_num, title, bullets, subtitle=None):
    """Generic slide with title + bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    # Accent bar
    add_accent_bar(slide, top=Inches(1.35))

    # Subtitle if provided
    y_start = Inches(1.7)
    if subtitle:
        txBox_sub = slide.shapes.add_textbox(Inches(0.8), y_start, Inches(11.7), Inches(0.5))
        tf_sub = txBox_sub.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        run_sub = p_sub.add_run()
        run_sub.text = subtitle
        run_sub.font.size = Pt(16)
        run_sub.font.color.rgb = LIGHT_GRAY
        run_sub.font.name = FONT_NAME
        y_start = Inches(2.3)

    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(1.0), y_start, Inches(11.3), Inches(4.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.space_after = Pt(12)
        p.space_before = Pt(4)

        # Teal bullet marker
        run_marker = p.add_run()
        run_marker.text = "\u25B8  "  # Small right-pointing triangle
        run_marker.font.size = Pt(18)
        run_marker.font.color.rgb = TEAL
        run_marker.font.name = FONT_NAME

        run_text = p.add_run()
        run_text.text = bullet
        run_text.font.size = Pt(18)
        run_text.font.color.rgb = WHITE
        run_text.font.name = FONT_NAME

    add_slide_number(slide, slide_num)
    return slide


def make_numbers_slide(prs, slide_num):
    """Slide 4: Platform by the Numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Platform by the Numbers"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    add_accent_bar(slide, top=Inches(1.35))

    # Number cards in a grid
    numbers = [
        ("25,054", "Evidence Claims"),
        ("515", "Hypotheses"),
        ("5,216", "Sources"),
        ("449", "Clinical Trials"),
        ("21 / 16", "Targets / Drugs"),
        ("29", "MCP Tools"),
        ("~190", "API Endpoints"),
    ]

    cols = 4
    card_w = Inches(2.6)
    card_h = Inches(1.8)
    margin_x = Inches(0.8)
    margin_y = Inches(2.0)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)

    for i, (num, label) in enumerate(numbers):
        col = i % cols
        row = i // cols
        left = margin_x + col * (card_w + gap_x)
        top = margin_y + row * (card_h + gap_y)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_LIGHTER
        card.line.color.rgb = DARK_BORDER
        card.line.width = Pt(1)

        # Number
        txNum = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), card_w - Inches(0.3), Inches(0.9))
        tfN = txNum.text_frame
        tfN.word_wrap = True
        pN = tfN.paragraphs[0]
        pN.alignment = PP_ALIGN.CENTER
        rN = pN.add_run()
        rN.text = num
        rN.font.size = Pt(36)
        rN.font.bold = True
        rN.font.color.rgb = TEAL
        rN.font.name = FONT_NAME

        # Label
        txLabel = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(1.05), card_w - Inches(0.3), Inches(0.6))
        tfL = txLabel.text_frame
        tfL.word_wrap = True
        pL = tfL.paragraphs[0]
        pL.alignment = PP_ALIGN.CENTER
        rL = pL.add_run()
        rL.text = label
        rL.font.size = Pt(14)
        rL.font.color.rgb = LIGHT_GRAY
        rL.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def make_pipeline_slide(prs, slide_num):
    """Slide 6: Evidence Pipeline (flow diagram)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "From Papers to Predictions"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    add_accent_bar(slide, top=Inches(1.35))

    # Flow steps
    steps = [
        "PubMed\n& Patents",
        "Claim\nExtraction",
        "Knowledge\nGraph",
        "Hypothesis\nGeneration",
        "GPU\nValidation",
    ]

    step_w = Inches(2.0)
    step_h = Inches(1.3)
    start_x = Inches(0.7)
    y = Inches(2.3)
    arrow_w = Inches(0.4)

    for i, step in enumerate(steps):
        left = start_x + i * (step_w + arrow_w)

        # Step box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, step_w, step_h)
        box.fill.solid()
        box.fill.fore_color.rgb = BG_LIGHTER
        box.line.color.rgb = TEAL
        box.line.width = Pt(2)

        # Text inside box
        tf_box = box.text_frame
        tf_box.word_wrap = True
        for j, line in enumerate(step.split("\n")):
            if j == 0:
                p_box = tf_box.paragraphs[0]
            else:
                p_box = tf_box.add_paragraph()
            p_box.alignment = PP_ALIGN.CENTER
            r_box = p_box.add_run()
            r_box.text = line
            r_box.font.size = Pt(16)
            r_box.font.bold = True
            r_box.font.color.rgb = WHITE
            r_box.font.name = FONT_NAME
        tf_box.paragraphs[0].space_before = Pt(14)

        # Arrow between steps
        if i < len(steps) - 1:
            arrow_left = left + step_w + Inches(0.05)
            arrow_y = y + step_h / 2 - Inches(0.15)
            txA = slide.shapes.add_textbox(arrow_left, arrow_y, Inches(0.35), Inches(0.4))
            tfA = txA.text_frame
            pA = tfA.paragraphs[0]
            pA.alignment = PP_ALIGN.CENTER
            rA = pA.add_run()
            rA.text = "\u25B6"  # Right arrow
            rA.font.size = Pt(20)
            rA.font.color.rgb = TEAL
            rA.font.name = FONT_NAME

    # Claim types below
    txClaims = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.5))
    tfC = txClaims.text_frame
    tfC.word_wrap = True
    pC = tfC.paragraphs[0]
    pC.alignment = PP_ALIGN.LEFT
    rC_label = pC.add_run()
    rC_label.text = "12 Claim Types:  "
    rC_label.font.size = Pt(14)
    rC_label.font.bold = True
    rC_label.font.color.rgb = TEAL
    rC_label.font.name = FONT_NAME

    claim_types = "gene_expression, protein_interaction, drug_efficacy, splicing_event, biomarker, pathway_mechanism, clinical_outcome, motor_function, neuroprotection, muscle_physiology, smn_regulation, therapeutic_target"
    rC_text = pC.add_run()
    rC_text.text = claim_types
    rC_text.font.size = Pt(12)
    rC_text.font.color.rgb = LIGHT_GRAY
    rC_text.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def make_targets_slide(prs, slide_num):
    """Slide 7: Key Targets table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "21 SMA Therapeutic Targets"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    add_accent_bar(slide, top=Inches(1.35))

    targets = [
        ("SMN1 / SMN2", "Primary SMA genes — survival of motor neuron"),
        ("PLS3", "Actin-bundling severity modifier"),
        ("STMN2", "Axonal growth regulator (TDP-43 connection)"),
        ("NCALD", "Calcium sensor — knockdown rescues SMA in mice"),
        ("UBA1", "Ubiquitin homeostasis — most druggable target"),
        ("CORO1C", "Actin remodeling — discovery target for 4-AP"),
    ]

    y_start = Inches(1.9)
    row_h = Inches(0.75)
    name_w = Inches(2.5)
    desc_w = Inches(9.0)

    for i, (name, desc) in enumerate(targets):
        y = y_start + i * row_h

        # Alternating row background
        if i % 2 == 0:
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.8), y - Inches(0.05),
                Inches(11.7), row_h - Inches(0.05)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = BG_LIGHTER
            row_bg.line.fill.background()

        # Target name
        txName = slide.shapes.add_textbox(Inches(1.0), y, name_w, row_h - Inches(0.1))
        tfN = txName.text_frame
        tfN.word_wrap = True
        tfN.vertical_anchor = MSO_ANCHOR.MIDDLE
        pN = tfN.paragraphs[0]
        rN = pN.add_run()
        rN.text = name
        rN.font.size = Pt(18)
        rN.font.bold = True
        rN.font.color.rgb = TEAL
        rN.font.name = FONT_NAME

        # Description
        txDesc = slide.shapes.add_textbox(Inches(3.8), y, desc_w, row_h - Inches(0.1))
        tfD = txDesc.text_frame
        tfD.word_wrap = True
        tfD.vertical_anchor = MSO_ANCHOR.MIDDLE
        pD = tfD.paragraphs[0]
        rD = pD.add_run()
        rD.text = desc
        rD.font.size = Pt(16)
        rD.font.color.rgb = WHITE
        rD.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def make_discovery_slide(prs, slide_num):
    """Slide 9: 4-AP Multi-Target Binding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "4-Aminopyridine: Multi-Target SMA Binder"
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    add_accent_bar(slide, top=Inches(1.35))

    # Key finding callout
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.0)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x22)  # Dark teal
    callout.line.color.rgb = TEAL
    callout.line.width = Pt(2)

    tf_call = callout.text_frame
    tf_call.word_wrap = True
    p_call = tf_call.paragraphs[0]
    p_call.alignment = PP_ALIGN.CENTER
    p_call.space_before = Pt(8)
    r_call = p_call.add_run()
    r_call.text = "KEY FINDING:  4-AP binds CORO1C with +0.251 confidence — strongest of ALL 378 docking pairs"
    r_call.font.size = Pt(18)
    r_call.font.bold = True
    r_call.font.color.rgb = TEAL
    r_call.font.name = FONT_NAME

    # Table of results
    results = [
        ("Target", "Confidence", "Note"),
        ("CORO1C", "+0.251", "STRONGEST"),
        ("NCALD", "-0.443", ""),
        ("SMN2", "-0.447", ""),
        ("SMN1", "-0.487", ""),
    ]

    col_widths = [Inches(3.0), Inches(2.5), Inches(3.0)]
    table_left = Inches(2.5)
    table_top = Inches(3.2)
    row_h = Inches(0.55)

    for i, (target, conf, note) in enumerate(results):
        y = table_top + i * row_h
        is_header = (i == 0)
        is_highlight = (i == 1)

        # Row background
        if is_header:
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, table_left, y,
                sum(w for w in col_widths), row_h
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = DARK_BORDER
            row_bg.line.fill.background()
        elif is_highlight:
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, table_left, y,
                sum(w for w in col_widths), row_h
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x22)
            row_bg.line.fill.background()
        elif i % 2 == 0:
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, table_left, y,
                sum(w for w in col_widths), row_h
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = BG_LIGHTER
            row_bg.line.fill.background()

        x = table_left
        for j, (text, width) in enumerate(zip([target, conf, note], col_widths)):
            txCell = slide.shapes.add_textbox(x + Inches(0.1), y, width - Inches(0.2), row_h)
            tfCell = txCell.text_frame
            tfCell.vertical_anchor = MSO_ANCHOR.MIDDLE
            pCell = tfCell.paragraphs[0]
            pCell.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            rCell = pCell.add_run()
            rCell.text = text
            rCell.font.name = FONT_NAME

            if is_header:
                rCell.font.size = Pt(14)
                rCell.font.bold = True
                rCell.font.color.rgb = LIGHT_GRAY
            elif is_highlight:
                rCell.font.size = Pt(16)
                rCell.font.bold = True
                rCell.font.color.rgb = TEAL
            else:
                rCell.font.size = Pt(15)
                rCell.font.color.rgb = WHITE

            x += width

    # Note
    txNote = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5))
    tfNote = txNote.text_frame
    tfNote.word_wrap = True
    pNote = tfNote.paragraphs[0]
    rNote = pNote.add_run()
    rNote.text = "Note: 4-AP is FDA-approved (multiple sclerosis) — previous SMA trial tested the WRONG mechanism"
    rNote.font.size = Pt(14)
    rNote.font.italic = True
    rNote.font.color.rgb = LIGHT_GRAY
    rNote.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def make_diffdock_slide(prs, slide_num):
    """Slide 10: DiffDock v2.2 Campaign."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "378 Docking Predictions via NVIDIA NIM"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    add_accent_bar(slide, top=Inches(1.35))

    # Key stats in cards
    stats = [
        ("54 x 7", "Compounds \u00d7 Targets"),
        ("378", "Blind Dockings"),
        ("3", "Positive-Confidence Hits"),
        ("24 min", "Total Runtime"),
        ("$0", "API Cost (Free Credits)"),
    ]

    card_w = Inches(2.2)
    card_h = Inches(1.5)
    start_x = Inches(0.6)
    y = Inches(1.9)
    gap = Inches(0.25)

    for i, (num, label) in enumerate(stats):
        left = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_LIGHTER
        card.line.color.rgb = DARK_BORDER
        card.line.width = Pt(1)

        txN = slide.shapes.add_textbox(left + Inches(0.1), y + Inches(0.2), card_w - Inches(0.2), Inches(0.7))
        tfN = txN.text_frame
        pN = tfN.paragraphs[0]
        pN.alignment = PP_ALIGN.CENTER
        rN = pN.add_run()
        rN.text = num
        rN.font.size = Pt(28)
        rN.font.bold = True
        rN.font.color.rgb = TEAL
        rN.font.name = FONT_NAME

        txL = slide.shapes.add_textbox(left + Inches(0.1), y + Inches(0.85), card_w - Inches(0.2), Inches(0.5))
        tfL = txL.text_frame
        tfL.word_wrap = True
        pL = tfL.paragraphs[0]
        pL.alignment = PP_ALIGN.CENTER
        rL = pL.add_run()
        rL.text = label
        rL.font.size = Pt(11)
        rL.font.color.rgb = LIGHT_GRAY
        rL.font.name = FONT_NAME

    # Additional context
    context_bullets = [
        "Only 3 of 378 pairs showed positive confidence (extremely selective)",
        "All 3 positive hits involve 4-Aminopyridine (4-AP)",
        "DiffDock v2.2 is 16% more accurate than v1 (NVIDIA GTC 2026)",
    ]

    txCtx = slide.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(11.3), Inches(2.5))
    tfCtx = txCtx.text_frame
    tfCtx.word_wrap = True

    for i, bullet in enumerate(context_bullets):
        if i == 0:
            p = tfCtx.paragraphs[0]
        else:
            p = tfCtx.add_paragraph()
        p.space_after = Pt(10)
        rm = p.add_run()
        rm.text = "\u25B8  "
        rm.font.size = Pt(16)
        rm.font.color.rgb = TEAL
        rm.font.name = FONT_NAME
        rt = p.add_run()
        rt.text = bullet
        rt.font.size = Pt(16)
        rt.font.color.rgb = WHITE
        rt.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def make_cross_paper_slide(prs, slide_num):
    """Slide 11: Cross-Paper Synthesis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Differentiator #1: Cross-Paper Synthesis'
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    add_accent_bar(slide, top=Inches(1.35))

    # Subtitle
    txSub = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.5))
    tfS = txSub.text_frame
    pS = tfS.paragraphs[0]
    rS = pS.add_run()
    rS.text = "Finding connections no single researcher would spot"
    rS.font.size = Pt(18)
    rS.font.italic = True
    rS.font.color.rgb = LIGHT_GRAY
    rS.font.name = FONT_NAME

    # Stats cards
    stats = [
        ("30", "Target\nCo-occurrence Pairs"),
        ("53", "Transitive Bridges\nAcross Papers"),
        ("6", "Claude Synthesis\nCards"),
    ]

    card_w = Inches(3.2)
    card_h = Inches(1.5)
    start_x = Inches(1.2)
    y = Inches(2.5)
    gap = Inches(0.5)

    for i, (num, label) in enumerate(stats):
        left = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_LIGHTER
        card.line.color.rgb = DARK_BORDER
        card.line.width = Pt(1)

        txN = slide.shapes.add_textbox(left, y + Inches(0.15), card_w, Inches(0.7))
        tfN = txN.text_frame
        pN = tfN.paragraphs[0]
        pN.alignment = PP_ALIGN.CENTER
        rN = pN.add_run()
        rN.text = num
        rN.font.size = Pt(36)
        rN.font.bold = True
        rN.font.color.rgb = TEAL
        rN.font.name = FONT_NAME

        txL = slide.shapes.add_textbox(left, y + Inches(0.85), card_w, Inches(0.55))
        tfL = txL.text_frame
        tfL.word_wrap = True
        pL = tfL.paragraphs[0]
        pL.alignment = PP_ALIGN.CENTER
        rL = pL.add_run()
        rL.text = label
        rL.font.size = Pt(13)
        rL.font.color.rgb = LIGHT_GRAY
        rL.font.name = FONT_NAME

    # Example bridge
    bridge_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.2)
    )
    bridge_box.fill.solid()
    bridge_box.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x22)
    bridge_box.line.color.rgb = TEAL
    bridge_box.line.width = Pt(1)

    txBridge = slide.shapes.add_textbox(Inches(1.8), Inches(4.65), Inches(9.7), Inches(1.0))
    tfB = txBridge.text_frame
    tfB.word_wrap = True
    pB1 = tfB.paragraphs[0]
    pB1.alignment = PP_ALIGN.CENTER
    rB1_label = pB1.add_run()
    rB1_label.text = "Example Bridge:  "
    rB1_label.font.size = Pt(14)
    rB1_label.font.bold = True
    rB1_label.font.color.rgb = TEAL
    rB1_label.font.name = FONT_NAME

    rB1_text = pB1.add_run()
    rB1_text.text = "SMN1  \u2192  SMN_PROTEIN  \u2192  NCALD  (calcium signaling bridge)"
    rB1_text.font.size = Pt(18)
    rB1_text.font.color.rgb = WHITE
    rB1_text.font.name = FONT_NAME

    pB2 = tfB.add_paragraph()
    pB2.alignment = PP_ALIGN.CENTER
    pB2.space_before = Pt(6)
    rB2 = pB2.add_run()
    rB2.text = "Two papers, two labs, one AI-discovered connection"
    rB2.font.size = Pt(13)
    rB2.font.italic = True
    rB2.font.color.rgb = LIGHT_GRAY
    rB2.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def make_cost_slide(prs, slide_num):
    """Slide 18: Cost Efficiency."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Drug Discovery for $12"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    add_accent_bar(slide, top=Inches(1.35))

    # Cost breakdown
    costs = [
        ("GPU Compute (Vast.ai A100)", "$0.78"),
        ("NIM API Docking (378 runs)", "$0.00"),
        ("Boltz-2 + ESM-2 (A100)", "~$0.05"),
        ("LLM Costs (Claude)", "~$11"),
        ("Total Session Cost", "~$12"),
    ]

    y_start = Inches(2.0)
    row_h = Inches(0.7)

    for i, (item, cost) in enumerate(costs):
        y = y_start + i * row_h
        is_total = (i == len(costs) - 1)

        if is_total:
            # Divider line
            div = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(2.0), y - Inches(0.1),
                Inches(9.3), Inches(0.03)
            )
            div.fill.solid()
            div.fill.fore_color.rgb = TEAL
            div.line.fill.background()
            y += Inches(0.1)

        # Item
        txItem = slide.shapes.add_textbox(Inches(2.0), y, Inches(7.0), row_h)
        tfI = txItem.text_frame
        tfI.vertical_anchor = MSO_ANCHOR.MIDDLE
        pI = tfI.paragraphs[0]
        rI = pI.add_run()
        rI.text = item
        rI.font.size = Pt(20) if is_total else Pt(18)
        rI.font.bold = is_total
        rI.font.color.rgb = WHITE
        rI.font.name = FONT_NAME

        # Cost
        txCost = slide.shapes.add_textbox(Inches(9.5), y, Inches(2.5), row_h)
        tfC = txCost.text_frame
        tfC.vertical_anchor = MSO_ANCHOR.MIDDLE
        pC = tfC.paragraphs[0]
        pC.alignment = PP_ALIGN.RIGHT
        rC = pC.add_run()
        rC.text = cost
        rC.font.size = Pt(22) if is_total else Pt(18)
        rC.font.bold = True
        rC.font.color.rgb = TEAL
        rC.font.name = FONT_NAME

    # Tagline
    txTag = slide.shapes.add_textbox(Inches(2.0), Inches(5.8), Inches(9.3), Inches(0.5))
    tfTag = txTag.text_frame
    pTag = tfTag.paragraphs[0]
    pTag.alignment = PP_ALIGN.CENTER
    rTag = pTag.add_run()
    rTag.text = "Democratizing computational drug discovery"
    rTag.font.size = Pt(18)
    rTag.font.italic = True
    rTag.font.color.rgb = LIGHT_GRAY
    rTag.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def make_smn1_ncald_slide(prs, slide_num):
    """Slide 15: SMN1-NCALD Bridge."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "AI Discovers: SMN1-NCALD Calcium Bridge"
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME

    add_accent_bar(slide, top=Inches(1.35))

    # Paper A box
    box_w = Inches(5.2)
    box_h = Inches(1.5)
    box_a = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), box_w, box_h
    )
    box_a.fill.solid()
    box_a.fill.fore_color.rgb = BG_LIGHTER
    box_a.line.color.rgb = DARK_BORDER
    box_a.line.width = Pt(1)

    txA = slide.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(4.8), box_h - Inches(0.3))
    tfA = txA.text_frame
    tfA.word_wrap = True
    pA1 = tfA.paragraphs[0]
    rA1 = pA1.add_run()
    rA1.text = "Paper A"
    rA1.font.size = Pt(16)
    rA1.font.bold = True
    rA1.font.color.rgb = TEAL
    rA1.font.name = FONT_NAME
    pA2 = tfA.add_paragraph()
    pA2.space_before = Pt(6)
    rA2 = pA2.add_run()
    rA2.text = "SMN1 loss \u2192 motor neuron degeneration"
    rA2.font.size = Pt(16)
    rA2.font.color.rgb = WHITE
    rA2.font.name = FONT_NAME

    # Paper B box
    box_b = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(2.0), box_w, box_h
    )
    box_b.fill.solid()
    box_b.fill.fore_color.rgb = BG_LIGHTER
    box_b.line.color.rgb = DARK_BORDER
    box_b.line.width = Pt(1)

    txB = slide.shapes.add_textbox(Inches(7.5), Inches(2.15), Inches(4.8), box_h - Inches(0.3))
    tfB = txB.text_frame
    tfB.word_wrap = True
    pB1 = tfB.paragraphs[0]
    rB1 = pB1.add_run()
    rB1.text = "Paper B"
    rB1.font.size = Pt(16)
    rB1.font.bold = True
    rB1.font.color.rgb = TEAL
    rB1.font.name = FONT_NAME
    pB2 = tfB.add_paragraph()
    pB2.space_before = Pt(6)
    rB2 = pB2.add_run()
    rB2.text = "NCALD-ASO treatment \u2192 rescues SMA mice"
    rB2.font.size = Pt(16)
    rB2.font.color.rgb = WHITE
    rB2.font.name = FONT_NAME

    # Arrow connecting
    arrow_mid = slide.shapes.add_textbox(Inches(5.8), Inches(2.3), Inches(1.7), Inches(0.8))
    tfArrow = arrow_mid.text_frame
    pArr = tfArrow.paragraphs[0]
    pArr.alignment = PP_ALIGN.CENTER
    rArr = pArr.add_run()
    rArr.text = "\u25C0 \u2500\u2500 \u25B6"
    rArr.font.size = Pt(22)
    rArr.font.color.rgb = TEAL
    rArr.font.name = FONT_NAME

    # AI connection box
    ai_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.0), Inches(8.3), Inches(1.8)
    )
    ai_box.fill.solid()
    ai_box.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x22)
    ai_box.line.color.rgb = TEAL
    ai_box.line.width = Pt(2)

    txAI = slide.shapes.add_textbox(Inches(2.8), Inches(4.15), Inches(7.7), Inches(1.5))
    tfAI = txAI.text_frame
    tfAI.word_wrap = True
    pAI1 = tfAI.paragraphs[0]
    pAI1.alignment = PP_ALIGN.CENTER
    rAI1 = pAI1.add_run()
    rAI1.text = "AI Cross-Paper Synthesis"
    rAI1.font.size = Pt(18)
    rAI1.font.bold = True
    rAI1.font.color.rgb = TEAL
    rAI1.font.name = FONT_NAME

    pAI2 = tfAI.add_paragraph()
    pAI2.alignment = PP_ALIGN.CENTER
    pAI2.space_before = Pt(8)
    rAI2 = pAI2.add_run()
    rAI2.text = "Connects both through calcium signaling pathway"
    rAI2.font.size = Pt(16)
    rAI2.font.color.rgb = WHITE
    rAI2.font.name = FONT_NAME

    pAI3 = tfAI.add_paragraph()
    pAI3.alignment = PP_ALIGN.CENTER
    pAI3.space_before = Pt(4)
    rAI3 = pAI3.add_run()
    rAI3.text = "Neither author saw this connection  |  Confidence: 0.52"
    rAI3.font.size = Pt(14)
    rAI3.font.italic = True
    rAI3.font.color.rgb = LIGHT_GRAY
    rAI3.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def make_thank_you_slide(prs, slide_num):
    """Slide 21: Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thank You"
    r.font.size = Pt(52)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT_NAME

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(3.4), Inches(3.0), Inches(0.04))
    div.fill.solid()
    div.fill.fore_color.rgb = TEAL
    div.line.fill.background()

    # Contact info
    contact_lines = [
        "sma-research.info",
        "christian@bryzant.com",
    ]

    txContact = slide.shapes.add_textbox(Inches(2.0), Inches(3.8), Inches(9.3), Inches(1.2))
    tfC = txContact.text_frame
    tfC.word_wrap = True

    for i, line in enumerate(contact_lines):
        if i == 0:
            p = tfC.paragraphs[0]
        else:
            p = tfC.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(20)
        r.font.color.rgb = TEAL
        r.font.name = FONT_NAME

    # Quote
    txQuote = slide.shapes.add_textbox(Inches(2.0), Inches(5.2), Inches(9.3), Inches(0.8))
    tfQ = txQuote.text_frame
    tfQ.word_wrap = True
    pQ = tfQ.paragraphs[0]
    pQ.alignment = PP_ALIGN.CENTER
    rQ = pQ.add_run()
    rQ.text = '"Built by a researcher with SMA. Biology-first. Evidence-only."'
    rQ.font.size = Pt(16)
    rQ.font.italic = True
    rQ.font.color.rgb = LIGHT_GRAY
    rQ.font.name = FONT_NAME

    add_slide_number(slide, slide_num)


def main():
    prs = Presentation()

    # Set widescreen 16:9
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1 — Title
    make_title_slide(prs, 1)

    # Slide 2 — The Problem
    make_bullet_slide(prs, 2, "The Unmet Need in SMA", [
        "SMA: #1 genetic cause of infant death \u2014 SMN1 gene loss",
        "3 approved therapies, but no cure \u2014 adults with long-standing disease respond poorly",
        "5,216+ papers published \u2014 no researcher can read them all",
        "Beyond-SMN targets (PLS3, NCALD, UBA1) largely unexplored computationally",
    ])

    # Slide 3 — Our Approach
    make_bullet_slide(prs, 3, "Evidence-First, AI-Driven, Open Source", [
        "Ingest ALL SMA literature (PubMed, patents, clinical trials)",
        "Extract structured claims using LLMs (Claude Haiku)",
        "Generate mechanistic hypotheses (Claude Sonnet 4.6)",
        "Validate with GPU-accelerated structural biology",
        "Everything traceable to source papers",
    ])

    # Slide 4 — Platform Scale (numbers)
    make_numbers_slide(prs, 4)

    # Slide 5 — Architecture
    make_bullet_slide(prs, 5, "Technical Architecture", [
        "FastAPI + PostgreSQL (asyncpg)",
        "Claude Sonnet 4.6 for reasoning, Haiku for extraction",
        "Vanilla JS frontend (no framework dependencies)",
        "Docker GPU toolkit (14.7GB, PyTorch + ESM-2 + OpenMM)",
        "NVIDIA NIM API integration (DiffDock v2.2, OpenFold3, GenMol)",
    ])

    # Slide 6 — Evidence Pipeline
    make_pipeline_slide(prs, 6)

    # Slide 7 — Key Targets
    make_targets_slide(prs, 7)

    # Slide 8 — GPU Infrastructure
    make_bullet_slide(prs, 8, "GPU-Accelerated Drug Discovery", [
        "Docker image: csiicf/sma-gpu-toolkit (14.7GB)",
        "Vast.ai on-demand A100 GPUs ($0.53/hr)",
        "Tools: DiffDock, Boltz-2, ESM-2, SpliceAI, OpenMM, Cas-OFFinder",
        "Total GPU cost: $0.78 for entire campaign",
        "First AI-driven virtual screening for SMA",
    ])

    # Slide 9 — Key Discovery: 4-AP
    make_discovery_slide(prs, 9)

    # Slide 10 — DiffDock Campaign
    make_diffdock_slide(prs, 10)

    # Slide 11 — Cross-Paper Synthesis
    make_cross_paper_slide(prs, 11)

    # Slide 12 — NVIDIA NIMs
    make_bullet_slide(prs, 12, "NVIDIA BioNeMo NIM Integration", [
        "DiffDock v2.2: Molecular docking (16% more accurate than v1)",
        "OpenFold3: Protein/RNA/ligand structure prediction",
        "GenMol: De novo molecule generation",
        "RNAPro: RNA 3D structure prediction",
        "All via cloud API \u2014 no local GPU needed",
    ])

    # Slide 13 — GenMol Results
    make_bullet_slide(prs, 13, "De Novo Molecule Design", [
        "Generated 20 novel 4-AP analogs via GenMol NIM",
        "Original 4-AP scaffold (+0.244) outperforms all analogs",
        "The aminopyridine core IS the optimal pharmacophore",
        "Implication: 4-AP structure is inherently suited for CORO1C binding",
    ])

    # Slide 14 — UBA1 Druggability
    make_bullet_slide(prs, 14, "UBA1: Most Druggable Target", [
        "5 compounds in top 25 docking hits",
        "CHEMBL1331875 (-0.089), CHEMBL1301743 (-0.100), ...",
        "UBA1 dysregulation is a known SMA feature",
        "Potential for ubiquitin-pathway targeted therapeutics",
        "Systematic druggability ranking across all 21 targets",
    ])

    # Slide 15 — SMN1-NCALD Bridge
    make_smn1_ncald_slide(prs, 15)

    # Slide 16 — MCP Server
    make_bullet_slide(prs, 16, "29 MCP Tools: Natural Language Access", [
        "Any AI assistant can query the entire SMA knowledge base",
        '"What is the evidence for PLS3 as a severity modifier?"',
        '"Show me compounds that bind SMN2"',
        '"What hypotheses connect NCALD to motor neuron survival?"',
        "Full programmatic access to 25,054 claims + 515 hypotheses",
    ], subtitle="Model Context Protocol \u2014 connecting AI assistants to live research data")

    # Slide 17 — Open Source
    make_bullet_slide(prs, 17, "Fully Open Source", [
        "AGPL-3.0 license \u2014 free to use, modify, and share",
        "GitHub: github.com/Bryzant-Labs/sma-platform",
        "HuggingFace: SMAResearch/sma-evidence-graph",
        "Docker Hub: csiicf/sma-gpu-toolkit",
        "All data traceable to source publications",
    ])

    # Slide 18 — Cost Efficiency
    make_cost_slide(prs, 18)

    # Slide 19 — Next Steps
    make_bullet_slide(prs, 19, "What's Next", [
        "100ns molecular dynamics simulation (OpenMM + CUDA)",
        "Experimental validation: SPR binding assay for 4-AP / CORO1C",
        "bioRxiv preprint: multi-target 4-AP binding",
        "Researcher outreach: Columbia, Cure SMA, EU labs",
        "Scale screening: 10K+ compounds",
    ])

    # Slide 20 — Limitations
    make_bullet_slide(prs, 20, "Honest Limitations", [
        "All results are computational predictions",
        "AlphaFold structures, not experimentally determined",
        "DiffDock confidence \u2260 binding affinity (Kd)",
        "Clinical trial failure of 4-AP in SMA is real",
        "Experimental validation required before any therapeutic claim",
    ])

    # Slide 21 — Thank You
    make_thank_you_slide(prs, 21)

    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "SMA_Research_Platform_EN.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
