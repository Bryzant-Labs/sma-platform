"""Generate SMA Research Platform presentation as PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0x0a, 0x0a, 0x0a)
CARD_BG = RGBColor(0x16, 0x16, 0x16)
WHITE = RGBColor(0xe8, 0xe8, 0xe8)
GRAY = RGBColor(0x88, 0x88, 0x88)
MUTED = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x00, 0xd4, 0xaa)
BLUE = RGBColor(0x4a, 0x9e, 0xff)
GOLD = RGBColor(0xff, 0xd7, 0x00)
RED = RGBColor(0xff, 0x47, 0x57)
PURPLE = RGBColor(0xa7, 0x8b, 0xfa)
LINE = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, font_name="Calibri", alignment=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def add_para(tf, text, font_size=14, color=GRAY, bold=False, space_before=0, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p


def add_rect(slide, left, top, width, height, fill_color=CARD_BG, border_color=LINE, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    return shape


def add_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_tag(slide, left, top, text, color=ACCENT, width=None):
    w = width or (len(text) * 0.11 + 0.3)
    rect = add_rect(slide, left, top, w, 0.32, fill_color=BG, border_color=color)
    txBox = add_text(slide, left, top, w, 0.32, text, font_size=9, color=color, font_name="Consolas", alignment=PP_ALIGN.CENTER)
    return rect


def slide_number(slide, num):
    add_text(slide, 0.5, 0.35, 0.8, 0.3, f"{num:02d}", font_size=10, color=MUTED, font_name="Consolas")


def section_label(slide, text, top=0.8):
    add_text(slide, 0.7, top, 4, 0.3, text.upper(), font_size=10, color=ACCENT, font_name="Consolas", bold=False)


# =============================================================
# SLIDE 1: TITLE
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s)

section_label(s, "Open-Source Forschungsplattform", top=1.5)
add_text(s, 0.7, 2.0, 10, 1.5, "SMA Research\nPlatform", font_size=54, color=WHITE, bold=True)
add_text(s, 0.7, 3.8, 9, 0.8, "Evidenz-basierte Wirkstoffforschung für Spinale Muskelatrophie —\nvon der automatisierten Literaturanalyse zur Hypothesenpriorisierung", font_size=16, color=GRAY)

add_tag(s, 0.7, 4.8, "EVIDENCE-FIRST", ACCENT, 1.5)
add_tag(s, 2.4, 4.8, "849 QUELLEN", BLUE, 1.3)
add_tag(s, 3.9, 4.8, "63 HYPOTHESEN", GOLD, 1.5)
add_tag(s, 5.6, 4.8, "10 TARGETS", GRAY, 1.3)

add_text(s, 0.7, 6.5, 6, 0.3, "Christian Fischer / Bryzant Labs  •  März 2026", font_size=11, color=MUTED)

# =============================================================
# SLIDE 2: MISSION
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 1)
section_label(s, "Mission")

add_text(s, 0.7, 1.3, 10, 0.8, "Warum diese Plattform?", font_size=36, color=WHITE, bold=True)
add_line(s, 0.7, 2.2, 1.5, ACCENT)

add_text(s, 0.7, 2.6, 9, 2.0,
    "SMA-Forschungsdaten sind über Hunderte von Datenbanken, Papers und Registern fragmentiert. "
    "Kein einzelner Forscher kann die gesamte Evidenz überblicken.\n\n"
    "Diese Plattform aggregiert, strukturiert und priorisiert die globale SMA-Evidenz "
    "automatisiert — damit kein vielversprechender Ansatz übersehen wird.",
    font_size=16, color=GRAY)

add_text(s, 0.7, 5.2, 9, 0.5,
    '\u201EJede unbeantwortete Frage ist ein Kind, das auf eine Therapie wartet.\u201C',
    font_size=14, color=MUTED, italic=True)

# =============================================================
# SLIDE 3: SMA KRANKHEITSBILD
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 2)
section_label(s, "Krankheitsbild")
add_text(s, 0.7, 1.3, 10, 0.6, "Spinale Muskelatrophie (SMA)", font_size=36, color=WHITE, bold=True)

cards = [
    ("Ursache", RED, "Homozygote Deletion/Mutation des SMN1-Gens auf Chromosom 5q13.\nVerlust des SMN-Proteins → Degeneration der Motoneuronen."),
    ("Modifier", GOLD, "SMN2-Kopienzahl bestimmt den Schweregrad.\nC-zu-T-Änderung in Exon 7 → ~90% Exon-Skipping → nur ~10% funktionelles Protein."),
    ("Prävalenz", BLUE, "~1:10.000 Geburten. Häufigste genetische Todesursache im Kindesalter.\nTrägerfrequenz: 1:35 – 1:50 in der Bevölkerung."),
    ("Therapien", PURPLE, "3 zugelassen: Nusinersen (ASO), Risdiplam (oral),\nOnasemnogene (Gentherapie). Heilung gibt es nicht."),
]

for i, (title, color, desc) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = 0.7 + col * 5.5
    y = 2.3 + row * 2.2
    add_rect(s, x, y, 5.0, 1.9, CARD_BG, LINE)
    add_text(s, x + 0.2, y + 0.15, 4.5, 0.35, title, font_size=14, color=color, bold=True)
    add_text(s, x + 0.2, y + 0.55, 4.5, 1.2, desc, font_size=12, color=GRAY)

# =============================================================
# SLIDE 4: ARCHITEKTUR
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 3)
section_label(s, "Architektur")
add_text(s, 0.7, 1.3, 10, 0.6, "Evidence-First Pipeline", font_size=36, color=WHITE, bold=True)

steps = [
    ("1. Ingestion", ACCENT, "PubMed, ClinicalTrials.gov, STRING, KEGG → 849 Quellen, 449 Studien"),
    ("2. Claim Extraction", BLUE, "LLM-basierte Extraktion strukturierter Aussagen aus Abstracts → 5.116 Claims"),
    ("3. Target Linking", GOLD, "Fuzzy-Matching + Drug-Target-Mapping → Claims werden Targets zugeordnet"),
    ("4. Priorisierung", PURPLE, "7-Dimensionen-Scoring + Hypothesen-Ranking → 63 Hypothesen in 3 Tiers"),
]

for i, (title, color, desc) in enumerate(steps):
    y = 2.2 + i * 1.15
    add_rect(s, 0.7, y, 7.0, 0.95, CARD_BG, LINE)
    # Colored left border
    add_rect(s, 0.7, y, 0.06, 0.95, color, color, Pt(0))
    add_text(s, 1.0, y + 0.08, 3, 0.3, title, font_size=12, color=color, bold=True, font_name="Consolas")
    add_text(s, 1.0, y + 0.42, 6.3, 0.45, desc, font_size=12, color=GRAY)

# Tech stack card
add_rect(s, 8.2, 2.2, 4.0, 4.6, CARD_BG, ACCENT)
add_text(s, 8.4, 2.35, 3.5, 0.3, "TECH STACK", font_size=10, color=ACCENT, font_name="Consolas", alignment=PP_ALIGN.CENTER)

tech = [
    ("Backend:", "Python / FastAPI"),
    ("Database:", "SQLite (asyncpg-kompatibel)"),
    ("LLM:", "Claude Haiku 4.5"),
    ("APIs:", "NCBI, STRING, KEGG"),
    ("Frontend:", "Vanilla JS"),
    ("Hosting:", "Moltbot Server"),
]
for i, (label, val) in enumerate(tech):
    y = 2.8 + i * 0.55
    add_text(s, 8.4, y, 1.5, 0.3, label, font_size=11, color=WHITE, bold=True)
    add_text(s, 9.9, y, 2.2, 0.3, val, font_size=11, color=GRAY)

# =============================================================
# SLIDE 5: ZAHLEN
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 4)
section_label(s, "Plattform-Status")
add_text(s, 0.7, 1.3, 10, 0.6, "Die Zahlen sprechen", font_size=36, color=WHITE, bold=True)

stats = [
    ("849", "PubMed-Quellen", ACCENT),
    ("5.116", "Extrahierte Claims", BLUE),
    ("63", "Hypothesen", GOLD),
    ("449", "Klinische Studien", PURPLE),
    ("10", "Mol. Targets", RED),
    ("7", "Wirkstoffe", GRAY),
]

for i, (num, label, color) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = 0.7 + col * 3.8
    y = 2.3 + row * 2.3
    add_rect(s, x, y, 3.3, 1.9, CARD_BG, LINE)
    add_text(s, x, y + 0.3, 3.3, 0.8, num, font_size=40, color=color, bold=True, font_name="Consolas", alignment=PP_ALIGN.CENTER)
    add_text(s, x, y + 1.2, 3.3, 0.4, label, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text(s, 0.7, 6.8, 10, 0.3, "Stand: März 2026 — automatisierte tägliche Aktualisierung via Cron um 06:00 UTC", font_size=10, color=MUTED)

# =============================================================
# SLIDE 6: TARGETS
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 5)
section_label(s, "Targets")
add_text(s, 0.7, 1.3, 10, 0.6, "10 molekulare Zielstrukturen", font_size=36, color=WHITE, bold=True)

# Table header
headers = [("Symbol", 1.8), ("Name", 4.0), ("Typ", 1.2), ("Claims", 1.2), ("Score", 1.0)]
x_start = 0.7
y_h = 2.3
x = x_start
for h, w in headers:
    add_text(s, x, y_h, w, 0.35, h, font_size=10, color=ACCENT, font_name="Consolas", bold=True)
    x += w
add_line(s, x_start, y_h + 0.35, 9.2, LINE)

targets = [
    ("SMN2", "Survival Motor Neuron 2", "Gen", "1.205", "89%", ACCENT),
    ("SMN1", "Survival Motor Neuron 1", "Gen", "934", "82%", BLUE),
    ("SMN_PROTEIN", "SMN Protein Complex", "Protein", "509", "60%", WHITE),
    ("NMJ_MATURATION", "Neuromuskuläre Junktion", "Pathway", "47", "38%", WHITE),
    ("MTOR_PATHWAY", "mTOR Signalweg", "Pathway", "44", "40%", WHITE),
    ("STMN2", "Stathmin-2", "Gen", "19", "38%", WHITE),
    ("UBA1", "Ubiquitin-Aktivierungsenzym", "Gen", "15", "40%", WHITE),
    ("PLS3", "Plastin 3", "Gen", "3", "31%", WHITE),
]

for i, (sym, name, typ, claims, score, sym_color) in enumerate(targets):
    y = 2.8 + i * 0.48
    x = x_start
    add_text(s, x, y, 1.8, 0.35, sym, font_size=11, color=sym_color, font_name="Consolas", bold=True)
    x += 1.8
    add_text(s, x, y, 4.0, 0.35, name, font_size=11, color=GRAY)
    x += 4.0
    add_text(s, x, y, 1.2, 0.35, typ, font_size=11, color=MUTED)
    x += 1.2
    add_text(s, x, y, 1.2, 0.35, claims, font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)
    x += 1.2
    sc_color = ACCENT if score.startswith("8") or score.startswith("9") else GOLD if score.startswith("6") else GRAY
    add_text(s, x, y, 1.0, 0.35, score, font_size=11, color=sc_color, font_name="Consolas", bold=True)

# =============================================================
# SLIDE 7: 7-DIMENSIONEN SCORING
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 6)
section_label(s, "Scoring")
add_text(s, 0.7, 1.3, 10, 0.6, "7-Dimensionen Target-Scoring", font_size=36, color=WHITE, bold=True)
add_text(s, 0.7, 2.0, 9, 0.4, "Beispiel: SMN2 — Composite Score: 89.0%", font_size=14, color=GRAY)

dims = [
    ("Evidence Volume", 100, ACCENT),
    ("Evidence Quality", 91, ACCENT),
    ("Bio. Coherence", 100, BLUE),
    ("Source Independence", 100, BLUE),
    ("Interventionability", 38, GOLD),
    ("Translational Readiness", 100, ACCENT),
    ("Network Centrality", 100, PURPLE),
]

for i, (label, pct, color) in enumerate(dims):
    y = 2.7 + i * 0.62
    add_text(s, 0.7, y, 2.8, 0.35, label, font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)
    # Track background
    add_rect(s, 3.7, y + 0.1, 6.5, 0.18, RGBColor(0x1a, 0x1a, 0x1a), RGBColor(0x1a, 0x1a, 0x1a), Pt(0))
    # Fill
    fill_w = 6.5 * (pct / 100)
    if fill_w > 0:
        add_rect(s, 3.7, y + 0.1, fill_w, 0.18, color, color, Pt(0))
    add_text(s, 10.4, y, 1.0, 0.35, f"{pct}%", font_size=11, color=GRAY, font_name="Consolas")

add_tag(s, 0.7, 7.0, "COMPOSITE SCORE", GOLD, 1.7)
add_text(s, 2.6, 6.9, 2, 0.45, "89.0%", font_size=28, color=ACCENT, bold=True, font_name="Consolas")

# =============================================================
# SLIDE 8: TOP 5 HYPOTHESEN
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 7)
section_label(s, "Phase 2 Ergebnis")
add_text(s, 0.7, 1.3, 10, 0.6, "Top 5 Hypothesen (Tier A)", font_size=36, color=WHITE, bold=True)

hyps = [
    ("#1", "SMN2", "64.0%", "SMN2-Kopienzahl als primärer Biomarker für Krankheitsschwere und Therapieansprechen"),
    ("#2", "SMN2", "61.5%", "SMN2-Splicing-Modulation durch Nusinersen/Risdiplam stellt funktionelles SMN-Protein wieder her"),
    ("#3", "SMN2", "60.8%", "ASO-basierte Neuroprotektion durch frühe SMN2-Korrektur maximiert Motoneuron-Erhalt"),
    ("#4", "SMN1", "60.1%", "SMN1-Kopienzahl + L-Arginin als prognostische Biomarker-Kombination"),
    ("#5", "SMN2", "60.0%", "U1 snRNP Splice-Site-Mechanismus als Ziel für Small-Molecule-Splicing-Modifier"),
]

for i, (rank, target, score, desc) in enumerate(hyps):
    y = 2.2 + i * 0.95
    add_rect(s, 0.7, y, 11.5, 0.8, CARD_BG, LINE)
    add_rect(s, 0.7, y, 0.06, 0.8, GOLD, GOLD, Pt(0))
    add_text(s, 1.0, y + 0.08, 0.6, 0.3, rank, font_size=11, color=GOLD, font_name="Consolas", bold=True)
    add_tag(s, 1.5, y + 0.1, target, BLUE if target == "SMN1" else ACCENT, 0.8)
    add_text(s, 10.8, y + 0.08, 1.2, 0.3, score, font_size=11, color=GOLD, font_name="Consolas", bold=True, alignment=PP_ALIGN.RIGHT)
    add_text(s, 1.0, y + 0.42, 11.0, 0.35, desc, font_size=12, color=GRAY)

# =============================================================
# SLIDE 9: THERAPIEN
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 8)
section_label(s, "Therapien")
add_text(s, 0.7, 1.3, 10, 0.6, "Zugelassene Wirkstoffe", font_size=36, color=WHITE, bold=True)

drugs = [
    ("Nusinersen", "SPINRAZA", ACCENT, "ASO — SMN2-Exon-7-Inklusion\nIntrathekal. Zugelassen 2016.", "458 Claims"),
    ("Risdiplam", "EVRYSDI", BLUE, "Small Molecule — SMN2-Splicing-Modifier\nOral. Zugelassen 2020.", "159 Claims"),
    ("Onasemnogene", "ZOLGENSMA", GOLD, "AAV9-Gentherapie — SMN1-Ersatz\nEinmalige IV-Infusion. Zugelassen 2019.", "184 Claims"),
]

for i, (name, brand, color, desc, claims) in enumerate(drugs):
    x = 0.7 + i * 3.9
    add_rect(s, x, 2.2, 3.5, 3.5, CARD_BG, LINE)
    # Top accent bar
    add_rect(s, x, 2.2, 3.5, 0.06, color, color, Pt(0))
    add_text(s, x + 0.2, 2.4, 3.0, 0.35, name, font_size=16, color=color, bold=True)
    add_text(s, x + 0.2, 2.85, 3.0, 0.25, brand, font_size=9, color=MUTED, font_name="Consolas")
    add_text(s, x + 0.2, 3.3, 3.0, 1.2, desc, font_size=12, color=GRAY)
    add_tag(s, x + 0.2, 5.2, claims, GRAY, 1.2)

add_text(s, 0.7, 6.2, 11, 0.3, "+ 4 weitere in der Pipeline: Apitegromab, Taldefgrobep, Reldesemtiv, Pyridostigmin", font_size=11, color=MUTED)

# =============================================================
# SLIDE 10: 3 PHASEN
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 9)
section_label(s, "Analyse-Toolchain")
add_text(s, 0.7, 1.3, 10, 0.6, "3 Phasen zur Wirkstoffentdeckung", font_size=36, color=WHITE, bold=True)

phases = [
    ("PHASE 1 ✓", ACCENT, "Evidenz & Hypothesen", "Ingestion → Normalisierung → Evidence Scoring →\nSignal Extraction → Hypothesenkarten", "ABGESCHLOSSEN", "Gate: 63/20-50 Hypothesen"),
    ("PHASE 2 ✓", GOLD, "Priorisierung", "Multi-Kriterien-Scoring → Tier A/B/C Ranking →\n5–15 High-Conviction Targets", "ABGESCHLOSSEN", "Output: 5 Tier A + 10 Tier B"),
    ("PHASE 3", MUTED, "Moleküldesign", "3A: Small Molecule (RDKit, Docking)\n3B: CRISPR/Prime Editing\n3C: AAV (Capsid, Tropismus)", "GEPLANT", "Benötigt: HuggingFace ML"),
]

for i, (phase, color, title, desc, status, note) in enumerate(phases):
    x = 0.7 + i * 3.9
    add_rect(s, x, 2.2, 3.5, 4.5, CARD_BG, color)
    add_rect(s, x, 2.2, 3.5, 0.06, color, color, Pt(0))
    add_text(s, x + 0.2, 2.4, 3.0, 0.25, phase, font_size=9, color=color, font_name="Consolas")
    add_text(s, x + 0.2, 2.75, 3.0, 0.35, title, font_size=14, color=WHITE, bold=True)
    add_text(s, x + 0.2, 3.3, 3.0, 1.5, desc, font_size=11, color=GRAY)
    add_tag(s, x + 0.2, 5.8, status, color, 1.5)
    add_text(s, x + 0.2, 6.2, 3.0, 0.3, note, font_size=9, color=MUTED)

# =============================================================
# SLIDE 11: MEILENSTEINE
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 10)
section_label(s, "Fortschritt")
add_text(s, 0.7, 1.3, 10, 0.6, "Erreichte Meilensteine", font_size=36, color=WHITE, bold=True)

milestones = [
    ("Phase 1A", ACCENT, "Datenbank & API", "10 Tabellen, FastAPI-Backend, SQLite, CORS-Security, Admin-Auth"),
    ("Phase 1B", ACCENT, "Ingestion-Adapter", "PubMed (49 Queries), ClinicalTrials.gov v2, STRING-DB, KEGG"),
    ("Phase 1C", ACCENT, "LLM Claim Extraction", "Claude Haiku extrahiert 5.116 Claims aus 849 Abstracts"),
    ("Phase 1D", ACCENT, "Netzwerk-Analyse", "STRING Protein-Interaktionen, Knowledge Graph, Network Centrality"),
    ("Phase 2A", GOLD, "Target-Scoring", "7-Dimensionen Priorisierung, Composite Scores, Tier-Labels"),
    ("Phase 2B", GOLD, "Hypothesen-Ranking", "63 Hypothesen in Tier A/B/C, 5-Dimensionen Scoring"),
    ("Phase 3", MUTED, "Moleküldesign", "CRISPR-Guides, AAV-Capsid-Evaluation (geplant)"),
]

for i, (phase, color, title, desc) in enumerate(milestones):
    y = 2.2 + i * 0.68
    # Dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.12), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # Line segment
    if i < len(milestones) - 1:
        line = add_rect(s, 1.04, y + 0.28, 0.04, 0.55, RGBColor(0x1a, 0x1a, 0x1a), RGBColor(0x1a, 0x1a, 0x1a), Pt(0))
        line.fill.solid()
        line.fill.fore_color.rgb = color
    add_text(s, 1.3, y, 1.3, 0.3, phase, font_size=9, color=color, font_name="Consolas")
    add_text(s, 2.7, y, 3.5, 0.3, title, font_size=13, color=WHITE, bold=True)
    add_text(s, 2.7, y + 0.3, 8, 0.3, desc, font_size=11, color=GRAY)

# =============================================================
# SLIDE 12: DATENQUELLEN
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 11)
section_label(s, "Datenquellen")
add_text(s, 0.7, 1.3, 10, 0.6, "Integrierte Datenbanken", font_size=36, color=WHITE, bold=True)

sources = [
    ("PubMed / NCBI", ACCENT, "849 Papers — 49 spezialisierte SMA-Queries,\nBiopython Entrez API, automatischer täglicher Pull"),
    ("ClinicalTrials.gov", BLUE, "449 Studien — REST v2 API,\nStatus-/Phasen-Tracking, Interventions-Mapping"),
    ("STRING-DB", GOLD, "Protein-Protein-Interaktionsnetzwerk —\n10 Edges für 7 SMA-Gene, Combined Scores"),
    ("GEO / NCBI", PURPLE, "7 kuratierte Omics-Datasets —\nTranskriptom, Proteom, Tiering nach Relevanz"),
]

for i, (name, color, desc) in enumerate(sources):
    col = i % 2
    row = i // 2
    x = 0.7 + col * 5.5
    y = 2.2 + row * 2.3
    add_rect(s, x, y, 5.0, 1.9, CARD_BG, LINE)
    add_text(s, x + 0.2, y + 0.2, 4.5, 0.35, name, font_size=14, color=color, bold=True)
    add_text(s, x + 0.2, y + 0.7, 4.5, 1.0, desc, font_size=12, color=GRAY)

add_text(s, 0.7, 6.8, 10, 0.3, "Weitere geplant: ChEMBL, DrugBank, UniProt, OpenTargets, AlphaFold", font_size=10, color=MUTED)

# =============================================================
# SLIDE 13: DAILY PIPELINE
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 12)
section_label(s, "Automatisierung")
add_text(s, 0.7, 1.3, 10, 0.6, "Tägliche Pipeline", font_size=36, color=WHITE, bold=True)
add_text(s, 0.7, 2.0, 9, 0.4, "Jeden Tag um 06:00 UTC läuft ein automatisierter Cron-Job:", font_size=14, color=GRAY)

cron_steps = [
    ("06:00", ACCENT, "PubMed-Pull — neue SMA-Papers der letzten 7 Tage"),
    ("06:01", BLUE, "ClinicalTrials.gov — alle SMA-Studien refreshen"),
    ("06:02", GOLD, "Claim Extraction — LLM-basierte Analyse neuer Abstracts"),
    ("06:05", PURPLE, "Claim Relinking — Drug-Target-Mapping + Fuzzy-Matching"),
    ("06:06", RED, "Score Refresh — Composite Scores aller Targets aktualisieren"),
]

for i, (time, color, desc) in enumerate(cron_steps):
    y = 2.7 + i * 0.85
    add_rect(s, 0.7, y, 11.0, 0.7, CARD_BG, LINE)
    add_text(s, 0.9, y + 0.18, 0.8, 0.3, time, font_size=11, color=color, font_name="Consolas", bold=True)
    add_text(s, 2.0, y + 0.18, 9.5, 0.3, desc, font_size=12, color=GRAY)

# =============================================================
# SLIDE 14: LINKS
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
slide_number(s, 13)
section_label(s, "Zugang")
add_text(s, 0.7, 1.3, 10, 0.6, "Links & Ressourcen", font_size=36, color=WHITE, bold=True)

links = [
    ("LIVE PLATTFORM", ACCENT, "https://sma-research.info", "Interaktives Dashboard mit allen Targets, Claims, Hypothesen und Scoring"),
    ("API ENDPOINT", BLUE, "https://sma-research.info/api/v2/", "RESTful API — /targets, /claims, /hypotheses/prioritized, /scores"),
    ("GITHUB", GRAY, "github.com/Bryzant-Labs/sma-platform", "Open Source — MIT License, Branch: master"),
]

for i, (label, color, url, desc) in enumerate(links):
    y = 2.3 + i * 1.5
    add_rect(s, 0.7, y, 8.0, 1.2, CARD_BG, LINE)
    add_rect(s, 0.7, y, 0.06, 1.2, color, color, Pt(0))
    add_text(s, 1.0, y + 0.1, 3, 0.25, label, font_size=9, color=color, font_name="Consolas")
    add_text(s, 1.0, y + 0.4, 7, 0.35, url, font_size=14, color=WHITE, bold=True)
    add_text(s, 1.0, y + 0.8, 7, 0.3, desc, font_size=11, color=GRAY)

# =============================================================
# SLIDE 15: CLOSING
# =============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_label(s, "Nächste Schritte", top=2.0)
add_text(s, 0, 2.5, 13.333, 0.8, "Bereit für Phase 3", font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(s, 2, 3.6, 9.333, 1.2,
    "63 priorisierte Hypothesen. 5 High-Conviction-Targets. Die Evidenzbasis steht.\n"
    "Als Nächstes: Computational Drug Discovery —\n"
    "Moleküldesign, CRISPR-Guides und AAV-Capsid-Evaluation.",
    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_tag(s, 4.0, 5.2, "SMA-RESEARCH.INFO", ACCENT, 2.0)
add_tag(s, 6.2, 5.2, "CHRISTIAN FISCHER", GRAY, 2.0)
add_tag(s, 8.4, 5.2, "OPEN SOURCE", BLUE, 1.5)

add_text(s, 0, 6.3, 13.333, 0.3, "Kontakt: christian@bryzant.com", font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)


# Save
output = "/mnt/c/Users/bryza/Dropbox/Christian fischer/SMA/sma-platform/docs/SMA_Research_Platform.pptx"
prs.save(output)
print(f"Saved: {output}")
